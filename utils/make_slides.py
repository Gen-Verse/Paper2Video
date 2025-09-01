from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import io
import aspose.slides as slides
import os
import subprocess

def make_slides(plan, ppt_file, image_file, new_file):
# create PPT
    prs = Presentation()

# adding ppt slide
    slide_layout = prs.slide_layouts[5] 
    slide = prs.slides.add_slide(slide_layout)

# adding title A
    title = slide.shapes.title
    title.text = plan["scenario"]
    img_path = image_file  
    left = Inches(1)
    top = Inches(1.5)
    height = Inches(3)
    slide.shapes.add_picture(img_path, left, top, height=height)

# adding C
    left = Inches(1)
    top = Inches(4.8)
    textbox = slide.shapes.add_textbox(left, top, Inches(8), Inches(1))
    text_frame = textbox.text_frame
    p = text_frame.add_paragraph()
    p.text = plan["prompt"]
    p.font.size = Pt(10)

# save document
    pptx_path = ppt_file
    prs.save(pptx_path)
    pptx_to_image(pptx_path,new_file)  # 假设已经将PPT转换为图片


def pptx_to_image(ppt_file, new_file):
    presentation = slides.Presentation(ppt_file)
    slide = presentation.slides[0]
    slide.get_thumbnail(1, 1).save(new_file)
def ppt_to_image2(ppt_file, output_dir):
    # make sure the path exists
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)
    # set the output path
    output_path = os.path.join(output_dir, "output.png")  # 输出的图片路径

    # use noconv with LibreOffice to convert PPT
    command = f"noconv -o {output_path} {ppt_file}"
    
    # run
    subprocess.run(command, shell=True, check=True)

    print(f"PPT is saved as image：{output_path}")


